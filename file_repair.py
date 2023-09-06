import os
from docx import Document
import openpyxl
import PyPDF2
from PIL import Image
from pptx import Presentation

def fix_word_document(file_path):
    document = Document(file_path)
    # Perform fixes on the document
    # Example: Replace all occurrences of 'old_word' with 'new_word'
    for paragraph in document.paragraphs:
        if 'old_word' in paragraph.text:
            paragraph.text = paragraph.text.replace('old_word', 'new_word')
    # Save the fixed document
    fixed_file_path = file_path.replace('all', 'repaired').replace('.docx', '_fixed.docx')
    document.save(fixed_file_path)


def fix_excel_spreadsheet(file_path):
    workbook = openpyxl.load_workbook(file_path)
    # Get the active sheet
    sheet = workbook.active
    # Perform fixes on the spreadsheet
    # Example: Replace all occurrences of 'old_value' with 'new_value'
    for row in sheet.iter_rows(values_only=True):
        for i, cell_value in enumerate(row):
            if cell_value == 'old_value':
                sheet.cell(row=row[0].row, column=i+1).value = 'new_value'
    # Save the fixed spreadsheet
    fixed_file_path = file_path.replace('all', 'repaired').replace('.xlsx', '_fixed.xlsx')
    workbook.save(fixed_file_path)


def fix_pdf_document(file_path):
    reader = PyPDF2.PdfReader(file_path)
    writer = PyPDF2.PdfWriter()
    # Perform fixes on the PDF
    # Example: Replace all occurrences of 'old_text' with 'new_text'
    for page in reader.pages:
        text = page.extract_text()
        fixed_text = text.replace('old_text', 'new_text')
        writer.add_page(page)
        writer.update_page_dict(page)
    # Save the fixed PDF
    fixed_file_path = file_path.replace('all', 'repaired').replace('.pdf', '_fixed.pdf')
    with open(fixed_file_path, 'wb') as output_file:
        writer.write(output_file)


def fix_image(file_path):
    image = Image.open(file_path)
    # Perform fixes on the image
    # Example: Rotate the image by 90 degrees
    fixed_image = image.rotate(90)
    # Save the fixed image
    fixed_file_path = file_path.replace('all', 'repaired').replace('.jpg', '_fixed.jpg')
    fixed_image.save(fixed_file_path)


def fix_powerpoint_presentation(file_path):
    presentation = Presentation(file_path)
    # Perform fixes on the presentation
    # Example: Replace all occurrences of 'old_text' with 'new_text'
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if 'old_text' in run.text:
                            run.text = run.text.replace('old_text', 'new_text')
    # Save the fixed presentation
    fixed_file_path = file_path.replace('all', 'repaired').replace('.pptx', '_fixed.pptx')
    presentation.save(fixed_file_path)


def fix_documents(file_path):
    if file_path.endswith('.docx'):
        fix_word_document(file_path)
    elif file_path.endswith('.xlsx'):
        fix_excel_spreadsheet(file_path)
    elif file_path.endswith('.pdf'):
        fix_pdf_document(file_path)
    elif file_path.endswith(('.jpg', '.jpeg', '.png')):
        fix_image(file_path)
    elif file_path.endswith('.pptx'):
        fix_powerpoint_presentation(file_path)
    else:
        print('Unsupported file format.')


# Traverse 'all' folder and fix the documents
all_dir = os.path.join(os.path.expanduser('~'), 'Documents', 'all')
repaired_dir = os.path.join(os.path.expanduser('~'), 'Documents', 'repaired')
unrepaired_dir = os.path.join(os.path.expanduser('~'), 'Documents', 'unrepaired')

for root, dirs, files in os.walk(all_dir):
    for file in files:
        file_path = os.path.join(root, file)
        if os.path.isfile(file_path):
            try:
                fix_documents(file_path)
                print(f"Document '{file}' repaired.")
                # Move the repaired document to the 'repaired' folder
                repaired_file_path = file_path.replace(all_dir, repaired_dir)
                os.rename(file_path, repaired_file_path)
            except:
                print(f"Failed to repair document '{file}'.")
                # Move the unrepaired document to the 'unrepaired' folder
                unrepaired_file_path = file_path.replace(all_dir, unrepaired_dir)
                os.rename(file_path, unrepaired_file_path)
        else:
            print(f"'{file_path}' is not a file.")

print("Repair process completed.")
